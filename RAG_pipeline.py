import json
import os
import base64     
import tempfile
from pathlib import Path
from typing import Dict, List, Optional

from dotenv import load_dotenv
from langchain_core.documents import Document
from langchain_core.messages import HumanMessage
from langchain_community.document_loaders import CSVLoader, PyPDFLoader
from langchain_community.vectorstores import Chroma
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from openai import OpenAI

try:
    from moviepy import VideoFileClip
    HAS_MOVIEPY = True
except ImportError:
    HAS_MOVIEPY = False
    print("Warning: 'moviepy' not found. Video processing will fail unless installed.")

# Load environment variables (.env for OPENAI_API_KEY)
load_dotenv()

# Load config from JSON so prompts/params are easy to tweak
CONFIG_PATH = Path("config.json")
if not CONFIG_PATH.exists():
    raise FileNotFoundError("config.json not found. Please create it before running the pipeline.")
with open(CONFIG_PATH, "r", encoding="utf-8") as f:
    CONFIG = json.load(f)

# Vector store location and supported file types
DATA_DIR = Path("Data")
VECTOR_DIR = DATA_DIR / "vector_store"
SUPPORTED_EXTS = {
    "pdf", "pptx", "csv", 
    "jpg", "jpeg", "png",       # Images
    "mp4", "mov", "avi", "mp3"  # Video/Audio
}
# Ensure vector store directory exists
def _ensure_dirs():
    """Ensures that the directory for the vector store exists."""
    VECTOR_DIR.mkdir(parents=True, exist_ok=True)

# Require OpenAI API key
def _require_openai_key():
    """
    Checks for the OpenAI API key in environment variables and raises an error if not found.

    This function serves as a guard to ensure that the necessary API key is available
    before any operations that require OpenAI services are attempted.

    Raises:
        ValueError: If the OPENAI_API_KEY environment variable is not set.
    """
    # Guard: ensure API key is available before any OpenAI call
    key = os.environ.get("OPENAI_API_KEY")
    if not key:
        print("Error: OPENAI_API_KEY is missing. Add it to your environment or .env file.")
        raise ValueError("OPENAI_API_KEY not set")

def _image_to_documents(file_path: Path) -> List[Document]:
    """
    Generates a semantic text description of an image file using GPT-4o (Vision).

    This function bridges the gap between visual data and text-based vector stores by 
    converting the visual content of an image into a detailed natural language description.

    The process involves:
    1. Reading the image file in binary mode.
    2. Converting the binary data to a Base64-encoded string to comport with OpenAI API requirements.
    3. Sending the encoded image to the `gpt-4o` model with a system prompt specifically designed 
       for retrieval databases (requesting details on visible text, charts, and key elements).
    4. Encapsulating the resulting description in a LangChain `Document` object.

    Args:
        file_path (Path): The file system path to the target image (.jpg, .jpeg, .png).

    Returns:
        List[Document]: A list containing a single Document object where `page_content` is the 
        model-generated description. Returns an empty list if an error occurs during API invocation 
        or file processing.

    Metadata:
        source (str): The absolute string path of the file.
        type (str): Fixed as "image".
    """   
    try:
        with open(file_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode("utf-8")
        
        vision_llm = ChatOpenAI(model="gpt-4o", temperature=0, max_tokens=1000)
        
        msg = HumanMessage(content=[
            {"type": "text", "text": "Describe this image in detail for a retrieval database. Include any visible text, charts, or key visual elements."},
            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{encoded_string}"}}
        ])
        
        response = vision_llm.invoke([msg])
        description = response.content
        
        return [Document(
            page_content=f"[Image Description of {file_path.name}]:\n{description}",
            metadata={
                "source": str(file_path),
                "type": "image"
            }
        )]
    except Exception as e:
        print(f"Error processing image {file_path}: {e}")
        return []
    

def _video_to_documents(file_path: Path) -> List[Document]:
    """
    Extracts and transcribes spoken audio from video or audio files using OpenAI Whisper.

    This function handles multimodal ingestion by isolating the audio track from a video file
     and converting speech to text.

    The process involves:
    1. Creating a temporary `.mp3` file to store the extracted audio track (if input is video).
    2. If input is video: extracting audio using `VideoFileClip`.
    3. If input is audio: using the file directly.
    4. Sending the audio binary to OpenAI's `whisper-1` model for transcription.
    5. Cleaning up any temporary files generated during the process.

    Args:
        file_path (Path): The file system path to the video (.mp4, .mov, .avi) or audio (.mp3) file.

    Returns:
        List[Document]: A list containing a single Document object where `page_content` is the 
        full text transcript. Returns an empty list if the transcript is empty or if an error occurs.

    Raises:
        ImportError: If the `moviepy` library is not installed.

    Metadata:
        source (str): The absolute string path of the file.
        type (str): Fixed as "video".
    """

    if not HAS_MOVIEPY:
        raise ImportError("moviepy is required for video processing. Run `pip install moviepy`.")

    try:
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as temp_audio:
            temp_audio_path = temp_audio.name
        
        if file_path.suffix.lower() in ['.mp3', '.wav']:
            temp_audio_path = str(file_path) 
            is_temp = False
        else:
            clip = VideoFileClip(str(file_path))
            clip.audio.write_audiofile(temp_audio_path, logger=None)
            is_temp = True

        client = OpenAI()
        with open(temp_audio_path, "rb") as audio_file:
            transcript = client.audio.transcriptions.create(
                model="whisper-1", 
                file=audio_file,
                prompt="Transcribe this audio clearly."
            )
        
        if is_temp and os.path.exists(temp_audio_path):
            os.remove(temp_audio_path)

        text_content = transcript.text
        if not text_content:
            return []

        return [Document(
            page_content=f"[Video/Audio Transcript of {file_path.name}]:\n{text_content}",
            metadata={
                "source": str(file_path),
                "type": "video"
            }
        )]

    except Exception as e:
        print(f"Error processing video {file_path}: {e}")
        if 'temp_audio_path' in locals() and os.path.exists(temp_audio_path) and is_temp:
            os.remove(temp_audio_path)
        return []

def _pptx_to_documents(file_path: Path) -> List[Document]:
    """
    Extracts raw text content from PowerPoint presentations (.pptx) without external service dependencies.

    This function iterates through the internal XML structure of a PowerPoint file to harvest text,
    maintaining the logical separation of content by creating one Document per slide.

    The process involves:
    1. Loading the presentation using `python-pptx`.
    2. Iterating through every slide in the deck.
    3. Within each slide, iterating through all shapes to find text-containing elements.
    4. Aggregating text within a slide into a single string.

    Args:
        file_path (Path): The file system path to the .pptx file.

    Returns:
        List[Document]: A list of Document objects, where each document represents one slide.

    Metadata:
        source (str): The absolute string path of the file.
        slide (int): The 1-based index of the slide.
        type (str): Fixed as "pptx".
    """
    prs = Presentation(str(file_path))
    docs: List[Document] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if not texts:
            continue
        content = "\n".join(texts)
        docs.append(
            Document(
                page_content=content,
                metadata={
                    "source": str(file_path),
                    "slide": idx,
                    "type": "pptx",
                },
            )
        )
    return docs

def _load_file_as_documents(file_path: Path) -> List[Document]:
    """
    The central dispatch function that routes files to their appropriate specific loader 
    based on file extension.

    Args:
        file_path (Path): The path to the file to be ingested.

    Returns:
        List[Document]: A list of processed Document objects ready for splitting/chunking.

    Raises:
        ValueError: If the file extension is not supported by any defined loader.
    """
    ext = file_path.suffix.lower()
    if ext == ".pdf":
        docs = PyPDFLoader(str(file_path)).load()
        for d in docs:
            d.metadata["type"] = "pdf"
        return docs
    if ext == ".csv":
        docs = CSVLoader(str(file_path), encoding="utf-8").load()
        for idx, d in enumerate(docs, start=1):
            d.metadata.setdefault("row", idx)
            d.metadata["type"] = "csv"
        return docs
    if ext == ".pptx":
        return _pptx_to_documents(file_path)
    
    if ext in {".jpg", ".jpeg", ".png"}:
        return _image_to_documents(file_path)
        
    if ext in {".mp4", ".mov", ".avi", ".mp3"}:
        return _video_to_documents(file_path)
    
    raise ValueError(f"Unsupported file type: {ext}")

def _split_documents(documents: List[Document]):
    """
    Splits a list of Documents into smaller, overlapping chunks suitable for vector embedding.

    The process relies on `RecursiveCharacterTextSplitter` and uses global `CONFIG` values for:
    - `CHUNK_SIZE`: The maximum size of a text block.
    - `CHUNK_OVERLAP`: The amount of text repeated between adjacent chunks to preserve context.

    Args:
        documents (List[Document]): The list of raw Document objects returned by the loaders.

    Returns:
        List[Document]: A new, longer list of Document objects where `page_content` has been 
        chunked to fit within specified size limits.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CONFIG["CHUNK_SIZE"],
        chunk_overlap=CONFIG["CHUNK_OVERLAP"],
        separators=["\n\n", "\n", " ", ""],
    )
    return splitter.split_documents(documents)


def _load_or_create_vector_store(embedding_model: OpenAIEmbeddings):
    """
    Loads an existing Chroma vector store from disk or returns None if it doesn't exist.

    This function checks for the presence of a persisted Chroma vector store in the
    directory specified by `VECTOR_DIR`. If the directory and a valid store are found,
    it initializes and returns a Chroma instance. Otherwise, it returns None, signaling
    that a new store needs to be created.

    Args:
        embedding_model (OpenAIEmbeddings): The embedding model to use for the vector store.

    Returns:
        Optional[Chroma]: A Chroma instance if the store exists, otherwise None.
    """
    # If the directory exists and has a Chroma collection, loading will reuse it
    if VECTOR_DIR.exists():
        return Chroma(persist_directory=str(VECTOR_DIR), embedding_function=embedding_model)
    return None


def _build_llm(model_override: Optional[str] = None, temperature_override: Optional[float] = None) -> ChatOpenAI:
    """
    Initializes and returns a ChatOpenAI client with configurable model and temperature.

    This function provides a centralized way to create an instance of the ChatOpenAI client.
    It uses default values from the global `CONFIG` but allows for overriding the model name
    and temperature on a per-call basis.

    Args:
        model_override (Optional[str]): If provided, this model name will be used instead of the one in `CONFIG`.
        temperature_override (Optional[float]): If provided, this temperature will be used instead of the one in `CONFIG`.

    Returns:
        ChatOpenAI: An instance of the ChatOpenAI client.
    """
    # Convenience to build an LLM client with optional overrides
    return ChatOpenAI(
        model=model_override or CONFIG["LLM_MODEL"],
        temperature=CONFIG["LLM_TEMPERATURE"] if temperature_override is None else temperature_override,
    )


def _generate_multi_queries(question: str, llm: ChatOpenAI) -> List[str]:
    """
    Generates alternative formulations of a user's question to improve retrieval recall.

    This function uses an LLM to rewrite a given question in multiple ways. These variants,
    along with the original question, can be used to query the vector store, increasing the
    likelihood of finding relevant documents even if the original phrasing is suboptimal.
    The number of variants is controlled by `MULTIQUERY_VARIANTS` in the `CONFIG`.

    Args:
        question (str): The original user question.
        llm (ChatOpenAI): The LLM client to use for generating variants.

    Returns:
        List[str]: A list of strings, with the original question as the first element,
                   followed by the generated variants.
    """
    # Generate multiple reformulations of the user query
    prompt = CONFIG["PROMPT_MULTIQUERY"].format(question=question)
    resp = llm.invoke(prompt)
    variants = [line.strip() for line in resp.content.splitlines() if line.strip()]
    if not variants:
        return [question]
    # Limit to configured count
    return [question] + variants[: CONFIG.get("MULTIQUERY_VARIANTS", 2)]


def _generate_hyde_text(query: str, llm: ChatOpenAI) -> str:
    """
    Generates a hypothetical document in response to a query (HyDE technique).

    Hypothetical Document Embeddings (HyDE) is a technique where a query is first used
    to generate a fictional document that answers it. The embedding of this hypothetical
    document is then used for the similarity search, which can sometimes yield more
    relevant results than embedding the raw query.

    Args:
        query (str): The user's query.
        llm (ChatOpenAI): The LLM client to use for generating the hypothetical document.

    Returns:
        str: The generated hypothetical document text.
    """
    # Generate a hypothetical answer to embed for retrieval (HyDE)
    prompt = CONFIG["PROMPT_HYDE"].format(question=query)
    resp = llm.invoke(prompt)
    return resp.content if hasattr(resp, "content") else str(resp)


def _dedup_docs(docs: List[Document], max_docs: int) -> List[Document]:
    """
    Deduplicates a list of Documents and caps the list at a maximum size.

    Deduplication is based on a composite key of the document's source, location
    (page, slide, or row), and the first 200 characters of its content. This helps
    ensure that the context provided to the LLM is both diverse and concise.

    Args:
        docs (List[Document]): The list of documents to deduplicate.
        max_docs (int): The maximum number of unique documents to return.

    Returns:
        List[Document]: A new list containing the unique documents, capped at `max_docs`.
    """
    # Drop duplicates by source+position+content prefix and cap the count
    seen = set()
    unique = []
    for doc in docs:
        key = (
            doc.metadata.get("source"),
            doc.metadata.get("page"),
            doc.metadata.get("slide"),
            doc.metadata.get("row"),
            doc.page_content[:200],
        )
        if key in seen:
            continue
        seen.add(key)
        unique.append(doc)
        if len(unique) >= max_docs:
            break
    return unique


def _normalize_source_meta(doc: Document) -> Dict:
    """
    Cleans and standardizes metadata from a Document for user-facing display.

    This function performs several transformations:
    - Extracts just the filename from the full source path.
    - Converts 0-based page numbers (from PyPDFLoader) to 1-based for readability.
    - Consolidates page, slide, and row numbers into a single "location" field.
    - Returns a dictionary with clean, presentation-ready metadata.

    Args:
        doc (Document): The Document object whose metadata needs to be normalized.

    Returns:
        Dict: A dictionary containing the normalized 'source', 'page', 'type', and 'location' fields.
    """
    # Normalize source name to filename and page/slide/row to human-friendly numbers
    raw_source = doc.metadata.get("source")
    source_name = Path(raw_source).name if raw_source else None

    page_val = doc.metadata.get("page")
    slide_val = doc.metadata.get("slide")
    row_val = doc.metadata.get("row")
    page_or_loc = page_val if page_val is not None else slide_val if slide_val is not None else row_val

    # If page is numeric, convert to int and +1 to show human page numbers (PyPDFLoader is 0-based)
    display_loc = page_or_loc
    if isinstance(page_or_loc, (int, float)) and page_or_loc >= 0:
        display_loc = int(page_or_loc) + 1
    else:
        # try parsing string digits
        try:
            num = int(str(page_or_loc))
            display_loc = num + 1
        except Exception:
            display_loc = page_or_loc

    return {
        "source": source_name,
        "page": display_loc if page_val is not None else None,
        "type": doc.metadata.get("type"),
        "location": display_loc,
    }


def add_file_to_rag(file_path: str, source_name: Optional[str] = None) -> Dict[str, int]:
    """
    Processes a single file, adds it to the RAG pipeline's vector store.

    This function orchestrates the entire ingestion process for a given file:
    1. Validates the file's existence and support for its extension.
    2. Routes the file to the appropriate loader (`_load_file_as_documents`).
    3. Splits the resulting documents into smaller chunks (`_split_documents`).
    4. Generates embeddings for each chunk using the configured OpenAI model.
    5. Adds the chunks to the Chroma vector store, creating the store if it doesn't exist.

    Args:
        file_path (str): The absolute or relative path to the file to be added.
        source_name (Optional[str]): An alternative name to store in the metadata. If not provided,
                                     the original filename is used.

    Returns:
        Dict[str, int]: A dictionary containing the number of chunks added and the new
                        total size of the collection.

    Raises:
        FileNotFoundError: If the specified `file_path` does not exist.
        ValueError: If the file type is not in the `SUPPORTED_EXTS` set.
    """
    # Reads the file (PDF/CSV/PPTX), chunks it, embeds with OpenAI, and updates the Chroma store.
    _ensure_dirs()
    _require_openai_key()

    path_obj = Path(file_path)
    if not path_obj.exists():
        print("Error: file not found.")
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path_obj.suffix.lower().lstrip(".")
    if ext not in SUPPORTED_EXTS:
        print("Error: unsupported file type.")
        raise ValueError(f"Unsupported file type: {ext}. Supported: {', '.join(SUPPORTED_EXTS)}")

    documents = _load_file_as_documents(path_obj)
    # Tag source as the original filename if provided (or use the incoming path name)
    src_name = Path(source_name).name if source_name else path_obj.name
    for d in documents:
        d.metadata["source"] = src_name
    chunks = _split_documents(documents)

    embedding = OpenAIEmbeddings(model=CONFIG["EMBEDDING_MODEL"])
    vector_store = _load_or_create_vector_store(embedding)

    if vector_store:
        vector_store.add_documents(chunks)
    else:
        vector_store = Chroma.from_documents(
            documents=chunks,
            embedding=embedding,
            persist_directory=str(VECTOR_DIR),
        )

    total = vector_store._collection.count() if hasattr(vector_store, "_collection") else -1

    return {"chunks_added": len(chunks), "collection_size": total}


def rag_query(
    question: str,
    k: Optional[int] = None,
    model: Optional[str] = None,
    temperature: Optional[float] = None,
    multiquery: Optional[bool] = None,
    hyde: Optional[bool] = None,
    mmr: Optional[bool] = None,
) -> Dict:
    """
    Executes a full Retrieval-Augmented Generation (RAG) query.

    This is the core function for asking questions. It orchestrates the following steps:
    1. Initializes the vector store and retriever.
    2. Optionally applies advanced retrieval strategies like Multi-Query and HyDE.
    3. Retrieves relevant document chunks from the vector store based on the query.
    4. Deduplicates and prepares the context.
    5. Constructs a final prompt including the retrieved context and the original question.
    6. Sends the prompt to the LLM to generate a final answer.
    7. Normalizes and returns the sources cited in the answer.

    Args:
        question (str): The user's question.
        k (Optional[int]): The number of documents to retrieve. Overrides `CONFIG`.
        model (Optional[str]): The LLM model to use. Overrides `CONFIG`.
        temperature (Optional[float]): The generation temperature. Overrides `CONFIG`.
        multiquery (Optional[bool]): Enables or disables Multi-Query retrieval. Overrides `CONFIG`.
        hyde (Optional[bool]): Enables or disables HyDE retrieval. Overrides `CONFIG`.
        mmr (Optional[bool]): Enables or disables Maximal Marginal Relevance (MMR) search. Overrides `CONFIG`.

    Returns:
        Dict: A dictionary containing the generated answer, a list of sources, and, if enabled,
              the queries and HyDE details used.

    Raises:
        ValueError: If the vector store has not been created yet.
    """
    # Runs a retrieval-augmented query using the persisted Chroma store.
    _require_openai_key()
    embedding = OpenAIEmbeddings(model=CONFIG["EMBEDDING_MODEL"])
    if not VECTOR_DIR.exists():
        print("Error: vector store not found. Ingest documents first.")
        raise ValueError("Vector store not found. Ingest documents first.")

    vector_store = Chroma(persist_directory=str(VECTOR_DIR), embedding_function=embedding)
    search_kwargs = {
        "k": k or CONFIG["RAG_SEARCH_K"],
    }
    if CONFIG.get("RAG_FETCH_K"):
        search_kwargs["fetch_k"] = CONFIG["RAG_FETCH_K"]
    if CONFIG.get("RAG_SCORE_THRESHOLD") is not None:
        search_kwargs["score_threshold"] = CONFIG["RAG_SCORE_THRESHOLD"]
    if CONFIG.get("RAG_FILTER_METADATA"):
        search_kwargs["filter"] = CONFIG["RAG_FILTER_METADATA"]

    search_type = "mmr" if mmr else CONFIG["RAG_SEARCH_TYPE"]
    retriever = vector_store.as_retriever(
        search_type=search_type,
        search_kwargs=search_kwargs,
    )

    llm = _build_llm(model_override=model, temperature_override=temperature)

    # Determine feature flags (request overrides config)
    use_multiquery = CONFIG.get("MULTIQUERY_ENABLED") if multiquery is None else multiquery
    use_hyde = CONFIG.get("HYDE_ENABLED") if hyde is None else hyde

    # Build query variants (multi-query) if enabled
    if use_multiquery:
        queries = _generate_multi_queries(question, llm)
    else:
        queries = [question]

    all_docs: List[Document] = []
    hyde_details = []
    for q in queries:
        query_text = q
        if use_hyde:
            query_text = _generate_hyde_text(q, llm)
            hyde_details.append({"original": q, "hyde": query_text})
        else:
            hyde_details.append({"original": q, "hyde": None})
        docs = retriever.invoke(query_text)
        all_docs.extend(docs)

    # Dedup and cap context docs
    context_docs = _dedup_docs(all_docs, max_docs=CONFIG.get("MAX_CONTEXT_DOCS", 8))

    # Format prompt
    context_text = "\n\n".join(doc.page_content for doc in context_docs)
    prompt = CONFIG["PROMPT_ANSWER"].format(context=context_text, question=question)
    response = llm.invoke(prompt)

    # Build a deduped source list for the response with normalized page numbers and filenames
    sources = []
    seen_sources = set()
    for doc in context_docs:
        normalized = _normalize_source_meta(doc)
        key = (normalized["source"], normalized["page"], normalized.get("type"))
        if key in seen_sources:
            continue
        seen_sources.add(key)
        sources.append(
            {
                "source": normalized["source"],
                "page": normalized["page"],
                "type": normalized.get("type"),
            }
        )
    answer_text = response.content if hasattr(response, "content") else str(response)
    result = {
        "answer": answer_text,
        "sources": sources,
    }
    if use_multiquery:
        result["queries"] = queries
    if use_hyde:
        result["hyde"] = hyde_details
    return result


def llm_only(
    question: str,
    model: Optional[str] = None,
    temperature: Optional[float] = None,
) -> str:
    """
    Queries the LLM directly without any retrieval augmentation.

    This function is useful for baseline comparisons to see how the LLM responds to a question
    using only its pre-trained knowledge, allowing for a direct assessment of the value
    added by the RAG pipeline.

    Args:
        question (str): The user's question.
        model (Optional[str]): The LLM model to use. Overrides `CONFIG`.
        temperature (Optional[float]): The generation temperature. Overrides `CONFIG`.

    Returns:
        str: The raw, non-augmented response from the LLM.
    """
    # Calls the LLM directly without retrieval for hallucination comparison.
    _require_openai_key()
    llm = ChatOpenAI(
        model=model or CONFIG["LLM_MODEL"],
        temperature=CONFIG["LLM_TEMPERATURE"] if temperature is None else temperature,
    )
    response = llm.invoke(question)
    return response.content if hasattr(response, "content") else response


def get_vector_store_stats() -> Dict[str, int]:
    """
    Retrieves statistics about the current state of the vector store.

    Checks if the vector store directory exists and, if so, returns the total
    number of documents (chunks) currently loaded in the Chroma collection.

    Returns:
        Dict[str, int]: A dictionary with the key 'documents_loaded' and its count.
                        Returns a count of 0 if the store does not exist.
    """
    # Returns statistics about the vector store.
    if not VECTOR_DIR.exists():
        return {"documents_loaded": 0}

    embedding = OpenAIEmbeddings(model=CONFIG["EMBEDDING_MODEL"])
    vector_store = Chroma(persist_directory=str(VECTOR_DIR), embedding_function=embedding)
    return {"documents_loaded": vector_store._collection.count()}


def list_source_files() -> List[str]:
    """
    Scans the vector store and returns a list of all unique source filenames.

    This function queries the metadata of all documents in the Chroma collection
    to compile a distinct, sorted list of the original files that have been ingested.

    Returns:
        List[str]: A sorted list of unique source filenames present in the vector store.
                   Returns an empty list if the store does not exist.
    """
    # Returns a list of unique source files in the vector store.
    if not VECTOR_DIR.exists():
        return []

    embedding = OpenAIEmbeddings(model=CONFIG["EMBEDDING_MODEL"])
    vector_store = Chroma(persist_directory=str(VECTOR_DIR), embedding_function=embedding)

    # Chroma's get() retrieves all documents and their metadata
    all_docs = vector_store.get(include=["metadatas"])

    # Extract unique source filenames from metadata
    sources = set()
    for metadata in all_docs.get("metadatas", []):
        if "source" in metadata:
            sources.add(metadata["source"])

    return sorted(list(sources))


def delete_source_file(filename: str) -> Dict[str, int]:
    """
    Deletes all document chunks from the vector store that originate from a specific file.

    This function finds all entries in the Chroma collection where the metadata 'source'
    matches the given filename and removes them. This is useful for managing the lifecycle
    of documents in the RAG system.

    Args:
        filename (str): The name of the source file to delete.

    Returns:
        Dict[str, int]: A dictionary containing the number of chunks that were deleted.

    Raises:
        ValueError: If the vector store does not exist or if the specified filename
                    is not found in the store.
    """
    # Deletes all chunks associated with a specific source file.
    if not VECTOR_DIR.exists():
        raise ValueError("Vector store not found.")

    embedding = OpenAIEmbeddings(model=CONFIG["EMBEDDING_MODEL"])
    vector_store = Chroma(persist_directory=str(VECTOR_DIR), embedding_function=embedding)

    # Find the IDs of all chunks with the given source filename to count them
    # and to ensure the file exists in the store.
    results = vector_store.get(where={"source": filename}, include=["metadatas"])
    ids_to_delete = results.get("ids", [])

    if not ids_to_delete:
        raise ValueError(f"File '{filename}' not found in vector store.")

    # Perform the deletion directly using a 'where' clause on the collection.
    # This is more atomic and reliable than deleting by IDs.
    vector_store._collection.delete(where={"source": filename})

    return {"chunks_deleted": len(ids_to_delete)}
